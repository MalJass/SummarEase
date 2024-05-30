import tkinter as tk
from tkinter import ttk
from tkinter import *
from tkinter import filedialog, messagebox
import requests
from tkinter.scrolledtext import *
import re
import nltk
import string
import numpy as np
import networkx as nx
from nltk.cluster.util import cosine_distance
import time
import PyPDF2
import docx
import pptx
from bs4 import BeautifulSoup
from urllib.request import urlopen




#time string for saving file
timestr =time.strftime("%Y%m%d-%H%M%S")
















#LOGIC FOR SUMMARY



# Function to generate summary
    
nltk.download('punkt')
nltk.download('stopwords')

stopwords = nltk.corpus.stopwords.words('english')

def preprocess(text):
    formatted_text = text.lower()
    tokens=[]
    for token in nltk.word_tokenize(formatted_text):
        tokens.append(token)
    tokens = [word for word in tokens if word not in stopwords and word not in string.punctuation]
    formatted_text = ' '.join(element for element in tokens)
    return formatted_text


def calculate_sentence_similarity(sentence1, sentence2):
  words1 = [word for word in nltk.word_tokenize(sentence1)]
  words2 = [word for word in nltk.word_tokenize(sentence2)]

  all_words = list(set(words1 + words2))

  vector1 = [0] * len(all_words)
  vector2 = [0] * len(all_words)

  for word in words1:
    vector1[all_words.index(word)] += 1
  for word in words2:
    vector2[all_words.index(word)] += 1
  
  return 1 - cosine_distance(vector1, vector2)

def calculate_similarity_matrix(sentences):
  similarity_matrix = np.zeros((len(sentences), len(sentences)))
  for i in range(len(sentences)):
    for j in range(len(sentences)):
      if i == j:
        continue
      similarity_matrix[i][j] = calculate_sentence_similarity(sentences[i], sentences[j])
  return similarity_matrix

def summarize(text, number_of_sentences, percentage = 0):
  original_s = [sentence for sentence in nltk.sent_tokenize(text)]
  formatted_s = [preprocess(original_s) for original_s in original_s]
  similarity_matrix = calculate_similarity_matrix(formatted_s)

  similarity_graph = nx.from_numpy_array(similarity_matrix)

  scores = nx.pagerank(similarity_graph)
  ordered_scores = sorted(((scores[i], score) for i, score in enumerate(original_s)), reverse=True)

  if percentage > 0:
    number_of_sentences = int(len(formatted_s) * percentage)

  best_sentences = []
  for sentence in range(number_of_sentences):
    best_sentences.append(ordered_scores[sentence][1])
    best_best ='\n'.join(best_sentences)
  return  best_best
 















#BUTTON FUNCTIONS





#Fuction of tab1 



# Function to display summary
def display_summary():
    try:
        original_s = entry.get("1.0", tk.END)
        best_best = summarize(original_s, 120 , 0.2)
        summary_text.delete("1.0",tk.END)
        summary_text.insert(tk.END, best_best)
    except Exception as e:
         messagebox.showerror("Error", f"An error occurred during summarization: {str(e)}")


# Function to reset text area
def clear_text_area():
    entry.delete(1.0, tk.END)

# Function to reset text area
def reset_text_area():
    summary_text.delete(1.0, tk.END)

#Function to save summary
def save():
    original_s =entry.get('1.0' ,tk.END)
    best_best = summarize(original_s,120 , 0.2)
    file_name ='your summary' + timestr + '.txt'
    with open(file_name , 'w') as f:
        f.write(best_best)
    summary_text.insert(tk.END, best_best)






#Fuction of file 


# Function to display summary
def display_summary_file():
    try:
        original_s = entry_file.get("1.0", tk.END)
        best_best = summarize(original_s, 120 , 0.2)
        summary_text_file.delete("1.0",tk.END)
        summary_text_file.insert(tk.END, best_best)
    except Exception as e:
         messagebox.showerror("Error", f"An error occurred during summarization: {str(e)}")


# Function to reset text area
def clear_text_file():
    entry_file.delete(1.0, tk.END)

# Function to reset text area
def reset_text_file():
    summary_text_file.delete(1.0, tk.END)

#Function to save summary
def save_file():
    original_s =entry_file.get('1.0' ,tk.END)
    best_best = summarize(original_s,120 , 0.2)
    file_name ='your summary' + timestr + '.txt'
    with open(file_name , 'w') as f:
        f.write(best_best)
    summary_text_file.insert(tk.END, best_best)


# Function to handle file upload
def upload_file():
    file_path = filedialog.askopenfilename()
    if file_path:
        try:
            if file_path.endswith('.pdf'):
                # Handle PDF files using PyPDF2
                with open(file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    text = ""
                    for page_num in range(len(pdf_reader.pages)):
                        text += pdf_reader.pages[page_num].extract_text() + "\n"
            elif file_path.endswith('.txt'):
                # Handle plain text files
                with open(file_path, 'r', encoding='utf-8') as file:
                    text = file.read()
            elif file_path.endswith('.docx'):
                # Handle DOCX files
                doc = docx.Document(file_path)
                paragraphs = [p.text for p in doc.paragraphs]
                text = '\n'.join(paragraphs)
            elif file_path.endswith('.pptx'):
                # Handle PPTX files
                prs = pptx.Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            else:
                # Unsupported file type
                messagebox.showerror("Error", "Unsupported file type. Please upload a PDF, TXT, DOCX, or PPTX file.")
                return

            # Clear the entry text area and insert the read text
            entry_file.delete('1.0', tk.END)
            entry_file.insert(tk.END, text)
        except Exception as e:
            messagebox.showerror("Error", f"Failed to read file: {e}")






#fuction of url


 
# Function to display summary
def display_summary_link():
    try:
        original_s = entry_link.get("1.0", tk.END)
        best_best = summarize(original_s, 120 , 0.2)
        summary_text_link.delete("1.0",tk.END)
        summary_text_link.insert(tk.END, best_best)
    except Exception as e:
         messagebox.showerror("Error", f"An error occurred during summarization: {str(e)}")


# Function to reset text area
def clear_text_link():
    entry_link.delete(1.0, tk.END)

# Function to reset text area
def reset_text_link():
    summary_text_link.delete(1.0, tk.END)

#Function to save summary
def save_link():
    original_s =entry_link.get('1.0' ,tk.END)
    best_best = summarize(original_s,120 , 0.2)
    file_name ='your summary' + timestr + '.txt'
    with open(file_name , 'w') as f:
        f.write(best_best)
    summary_text_link.insert(tk.END, best_best)


# Function to handle link input
def process_link():
    try:
        link = str(link_entry.get('1.0', tk.END))
        if link:
            page=urlopen(link)
            soup=BeautifulSoup(page,'lxml')
            text=' '.join(map(lambda p: p.text, soup.find_all('p')))
            entry_link.insert(tk.END , text)
        else:
            messagebox.showerror("Error", "Please enter a valid link.")
    except Exception as e:
        messagebox.showerror("Error", f"An error occurred: {str(e)}")




















#GUI



#window 
window= tk.Tk()
window.title("SummarEase")
window.geometry("800x500")

#style
style = ttk.Style(window)
style.configure('righttab.TNotebook' , tabposition='en')
#tabs
tab_control = ttk.Notebook(window,style ='righttab.TNotebook')

tab1= ttk.Frame(tab_control)
tab2= ttk.Frame(tab_control)
tab3= ttk.Frame(tab_control)
tab4= ttk.Frame(tab_control)

#add tabs to notebook
tab_control.add(tab1, text =f'{"Home":^16s}')
tab_control.add(tab2, text =f'{"File":^20s}')
tab_control.add(tab3, text =f'{"URL":^18s}')
tab_control.add(tab4, text =f'{"About":^17s}')










#HOME


#labels home

enter_text_label = tk.Label(tab1, text="Enter Your Text Here:",padx=5, pady=5)
enter_text_label.grid(row=1, column=0, sticky='nsw')

summary_label = tk.Label(tab1, text="Summary:", padx=5, pady=5)
summary_label.grid(row=1, column=2,sticky='nsw')

#buttons home

reset_button = tk.Button(tab1, text="Reset", command=reset_text_area, width=10, height=1, bg='#FF6347',fg ='#fff')
reset_button.grid(row=0, column=0,padx=10, pady=10,sticky='sne')

summary_button = tk.Button(tab1, text="Summarize", command=display_summary, width=10, height=1,bg='#FFA500')
summary_button.grid(row=0, column=3,padx=10, pady=10,sticky='sn')

clear_button = tk.Button(tab1, text="clear", command=clear_text_area, width=10, height=1, bg='#FF8C00',fg ='#fff')
clear_button.grid(row=0, column=1,padx=10, pady=10,sticky='sn')

save_button = tk.Button(tab1, text="Save", command=save, width=10, height=1, bg ='#FF7F50')
save_button.grid(row=0, column=2,padx=10, pady=10,sticky='sne')

#text area home

entry = ScrolledText(tab1, height=10, width=20)
entry.grid(row=2, column=0,columnspan=2,padx=5,pady=5,sticky='nsew')

summary_text = ScrolledText(tab1, wrap=tk.WORD, height=10, width=20)
summary_text.grid(row=2, column=2, columnspan=2,sticky='nsew')










#File


#labels file
original_text_label = tk.Label(tab2, text="Original:",padx=5, pady=5)
original_text_label.grid(row=2, column=0,sticky='nsw')

enter_file_label = tk.Label(tab2, text="Select the file to upload:",padx=5, pady=5)
enter_file_label.grid(row=0, column=0,columnspan=2,sticky='nse')

summary_label = tk.Label(tab2, text="Summary:", padx=5, pady=5)
summary_label.grid(row=2, column=2,sticky='nsw')

#buttons file

browse_button = tk.Button(tab2, text="Browse", command=upload_file, width=10, height=1, bg='#FF00FF',fg ='#fff')
browse_button.grid(row=0, column=2,columnspan=2,padx=10, pady=10,sticky='nw')


reset_button_file = tk.Button(tab2, text="Reset", command=reset_text_file, width=10, height=1, bg='#FF6347',fg ='#fff')
reset_button_file.grid(row=1, column=0,padx=10, pady=10,sticky='ns')

summary_button_file = tk.Button(tab2, text="Summarize", command=display_summary_file, width=10, height=1,bg='#FFA500')
summary_button_file.grid(row=1, column=3,padx=10, pady=10,sticky='ns')

clear_button_file = tk.Button(tab2, text="clear", command=clear_text_file, width=10, height=1, bg='#FF8C00',fg ='#fff')
clear_button_file.grid(row=1, column=1,padx=10, pady=10,sticky='nsw')

save_button_file = tk.Button(tab2, text="Save", command=save_file, width=10, height=1, bg ='#FF7F50')
save_button_file.grid(row=1, column=2,padx=10, pady=10,sticky='nse')

#text area file

entry_file = ScrolledText(tab2, wrap="word",height=10, width=20)
entry_file.grid(row=3, column=0,columnspan=2,padx=5,pady=5,sticky='nsew')

summary_text_file = ScrolledText(tab2, wrap=tk.WORD, height=10, width=20)
summary_text_file.grid(row=3, column=2, columnspan=3,sticky='nsew')








#URL

#labels url
original_text_label = tk.Label(tab3, text="Original:",padx=5, pady=5)
original_text_label.grid(row=2, column=0,sticky='nsw')

enter_url_label = tk.Label(tab3, text="Enter URL Here:",padx=5, pady=5)
enter_url_label.grid(row=0, column=0,sticky='nse')

summary_label = tk.Label(tab3, text="Summary:", padx=5, pady=5)
summary_label.grid(row=2, column=2,sticky='nsw')

#buttons url

process_button = tk.Button(tab3, text="Process", command=process_link)
process_button.grid(row=0, column=3,sticky='nsw')

reset_button_link = tk.Button(tab3, text="Reset", command=reset_text_link, width=10, height=1, bg='#FF6347',fg ='#fff')
reset_button_link.grid(row=1, column=0,padx=10, pady=10,sticky='ns')

summary_button_link = tk.Button(tab3, text="Summarize", command=display_summary_link, width=10, height=1,bg='#FFA500')
summary_button_link.grid(row=1, column=3,padx=10, pady=10,sticky='ns')

clear_button_link = tk.Button(tab3, text="clear", command=clear_text_link, width=10, height=1, bg='#FF8C00',fg ='#fff')
clear_button_link.grid(row=1, column=1,padx=10, pady=10,sticky='ns')

save_button_link = tk.Button(tab3, text="Save", command=save_link, width=10, height=1, bg ='#FF7F50')
save_button_link.grid(row=1, column=2,padx=10, pady=10,sticky='ns')

#text area url

link_entry = Text(tab3, height=1, width=20)
link_entry.grid(row=0, column=1,columnspan=2,padx=10, pady=10,sticky='nsew')

entry_link = ScrolledText(tab3, wrap="word",height=10, width=20)
entry_link.grid(row=3, column=0,columnspan=2,padx=5,pady=5,sticky='nsew')

summary_text_link = ScrolledText(tab3, wrap=tk.WORD, height=10, width=20)
summary_text_link.grid(row=3, column=2, columnspan=3,sticky='nsew')




#ABOUT

# Create the "About" section in Tab 4
about_label = tk.Label(tab4, text="About", font=("Helvetica", 16, "bold"))
about_label.grid(row=0, column=0, padx=10, pady=10, sticky="w")

# Information about the application
about_text = """
Text Summarizer v1.0

This application allows you to summarize text, files, or URLs using advanced natural language processing techniques.

Credits:
- NLTK: Natural Language Toolkit for text processing
- PyPDF2: Library for working with PDF files
- BeautifulSoup: Library for web scraping
- PyInstaller: Used for packaging the application as an executable

For support or feedback, please visit:
https://github.com/yourusername/text_summarizer

"""

about_info = tk.Label(tab4, text=about_text, justify="left")
about_info.grid(row=1, column=0, padx=10, pady=10, sticky="w")





tab_control.pack(expand=1,fill='both')



for j in range(4):  # Adjust the range according to the number of columns
    tab1.columnconfigure(j, weight=1)
    tab2.columnconfigure(j, weight=1)
    tab3.columnconfigure(j, weight=1)
    tab4.columnconfigure(j, weight=1)


tab1.rowconfigure(2, weight=1)
tab3.rowconfigure(3, weight=1)
tab2.rowconfigure(3, weight=1)
window.mainloop()
    

